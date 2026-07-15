from tools.document_parser import extract_text_from_file


def test_extract_text_markdown(tmp_path):
    md_file = tmp_path / "test.md"
    md_file.write_text("# Title\nThis is a markdown file.", encoding="utf-8")

    extracted = extract_text_from_file(md_file)
    assert "# Title" in extracted
    assert "markdown file" in extracted


def test_extract_text_docx(tmp_path):
    import docx

    docx_file = tmp_path / "test.docx"
    doc = docx.Document()
    doc.add_paragraph("Hello from Word document.")
    # Add a table
    table = doc.add_table(rows=1, cols=2)
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = "Header 1"
    hdr_cells[1].text = "Header 2"
    doc.save(docx_file)

    extracted = extract_text_from_file(docx_file)
    assert "Hello from Word document" in extracted
    assert "Header 1 | Header 2" in extracted


def test_extract_text_pptx(tmp_path):
    from pptx import Presentation

    pptx_file = tmp_path / "test.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    assert title is not None
    title.text = "PowerPoint Slide Title"
    prs.save(pptx_file)

    extracted = extract_text_from_file(pptx_file)
    assert "PowerPoint Slide Title" in extracted


def test_extract_text_xlsx(tmp_path):
    import openpyxl

    xlsx_file = tmp_path / "test.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    assert ws is not None
    ws.title = "TestSheet"
    ws.append(["Name", "Value"])
    ws.append(["Item 1", 100])
    wb.save(xlsx_file)

    extracted = extract_text_from_file(xlsx_file)
    assert "Sheet: TestSheet" in extracted
    assert "Name | Value" in extracted
    assert "Item 1 | 100" in extracted


def test_extract_text_pdf(tmp_path):
    from fpdf import FPDF
    from fpdf.enums import XPos, YPos

    pdf_file = tmp_path / "test.pdf"
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("helvetica", size=12)
    pdf.cell(
        200,
        10,
        text="Hello PDF File Content",
        new_x=XPos.LMARGIN,
        new_y=YPos.NEXT,
    )
    pdf.output(str(pdf_file))

    extracted = extract_text_from_file(pdf_file)
    assert "Hello PDF File Content" in extracted


def test_unsupported_and_truncation(tmp_path):
    # Test unsupported binary formats
    legacy_file = tmp_path / "test.doc"
    legacy_file.write_bytes(b"\xd0\xcf\x11\xe0")
    assert "Unsupported legacy binary file format" in extract_text_from_file(
        legacy_file
    )

    # Test truncation
    large_file = tmp_path / "large.txt"
    large_file.write_text("A" * 100, encoding="utf-8")
    extracted = extract_text_from_file(large_file, max_chars=10)
    assert "AAAAAA" in extracted
    assert "[TRUNCATED" in extracted
