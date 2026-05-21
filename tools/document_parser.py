import os
from pathlib import Path
import logging

logger = logging.getLogger(__name__)

def _get_int_env(name: str, default: int, min_value: int = 1) -> int:
    raw = os.getenv(name)
    if not raw:
        return default
    try:
        value = int(raw)
    except ValueError:
        return default
    if value < min_value:
        return default
    return value

MAX_READ_CHARS = _get_int_env("MAX_FILE_READ_CHARS", 120000)

def extract_text_from_file(file_path: Path, max_chars: int = MAX_READ_CHARS) -> str:
    """Extract text from PDF, DOCX, PPTX, XLSX, MD, or text files."""
    if isinstance(file_path, str):
        file_path = Path(file_path)

    if not file_path.is_file():
        return f"Error: File '{file_path}' does not exist or is not a file."

    suffix = file_path.suffix.lower()
    content = ""

    try:
        if suffix == ".pdf":
            # Native support for PDF via PyMuPDF (fitz)
            import fitz
            doc = fitz.open(file_path)
            pages_text = []
            for i, page in enumerate(doc, 1):
                page_text = page.get_text()
                if page_text.strip():
                    pages_text.append(f"--- Page {i} ---\n{page_text}")
            content = "\n\n".join(pages_text)

        elif suffix == ".docx":
            # Native support for Word documents
            import docx
            doc = docx.Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            
            # Format tables
            tables_text = []
            for t_idx, table in enumerate(doc.tables, 1):
                table_lines = [f"\n[Table {t_idx}]"]
                for row in table.rows:
                    row_cells = [cell.text.strip() for cell in row.cells]
                    table_lines.append(" | ".join(row_cells))
                tables_text.append("\n".join(table_lines))
                
            content = "\n\n".join(paragraphs)
            if tables_text:
                content += "\n\n" + "\n\n".join(tables_text)

        elif suffix == ".pptx":
            # Native support for PowerPoint presentations
            from pptx import Presentation
            prs = Presentation(file_path)
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                slide_lines = [f"--- Slide {i} ---"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_lines.append(shape.text.strip())
                slides_text.append("\n".join(slide_lines))
            content = "\n\n".join(slides_text)

        elif suffix == ".xlsx":
            # Native support for Excel spreadsheets
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            sheets_text = []
            for sheet in wb.worksheets:
                sheet_lines = [f"--- Sheet: {sheet.title} ---"]
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        row_str = " | ".join(str(cell) if cell is not None else "" for cell in row)
                        sheet_lines.append(row_str)
                sheets_text.append("\n".join(sheet_lines))
            content = "\n\n".join(sheets_text)

        elif suffix in (".md", ".txt", ".json", ".xml", ".yaml", ".yml", ".py", ".js", ".ts", ".html", ".css", ".csv"):
            # Plain text files
            content = file_path.read_text(encoding="utf-8", errors="replace")

        elif suffix in (".doc", ".ppt", ".xls"):
            return f"[Unsupported legacy binary file format: {suffix}. Please save/convert to the newer OpenXML format (e.g. .docx, .pptx, .xlsx) for native parsing.]"

        else:
            # Try to decode as utf-8, fallback to ASCII/error replace
            try:
                content = file_path.read_text(encoding="utf-8")
            except Exception:
                return f"[Unsupported binary file format: {suffix}]"

    except Exception as e:
        logger.exception(f"Error parsing document {file_path}")
        return f"Error parsing document '{file_path.name}': {e}"

    if len(content) > max_chars:
        return content[:max_chars] + f"\n\n[TRUNCATED: showing first {max_chars} characters of {len(content)} total]"

    return content
